# -*- coding: utf-8 -*-
"""OC2026 学科紹介スライド ビルダー（大きな文字・少ない情報量を強制する設計）"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree

# ---------------------------------------------------------------- design tokens
FONT = "Meiryo"          # Windows 全機種にある高可読ゴシック。BIZ UDPゴシックに一括置換も可
SLIDE_W = 13.3333
SLIDE_H = 7.5

# 会場投影を前提にコントラスト比を検証済み（本文7:1以上／24pt以上の太字で4.5:1以上）
C = {
    "ink":     "12253F",   # 濃紺（本文）           白地に 14.0:1
    "navy":    "0E2A4F",   # 背景の紺
    "navy_dk": "081A33",
    "brand":   "1160C4",   # ブランドブルー（面用）
    "brand_dk":"0B4C9E",   # ブランドブルー（文字用） 白地に 7.4:1
    "brand_lt":"E7F0FB",
    "gold":    "F2B01E",   # ゴールド（面だけに使う。白地の文字には使わない）
    "gold_lt": "FFF4DA",
    "gold_tx": "5C4204",   # 金の面の上の文字        8.6:1
    "coral":   "B34A33",   # 数字の強調（大きな文字のみ）5.3:1
    "teal":    "0A6156",   # 白抜き 7.4:1
    "teal_lt": "E1F3F1",
    "gray":    "41505F",   # 補足テキスト           白地に 8.3:1
    "gray_lt": "AFB9C5",   # 罫線専用。文字には使わない
    "hint":    "41505F",   # プレースホルダの指示文
    "line":    "D8DEE7",
    "bg":      "F5F7FA",
    "white":   "FFFFFF",
}

# タイプスケール（pt）: 会場後方・保護者でも読める下限を守る
T = {
    "kicker": 20,
    "title":  44,
    "title_s":38,
    "lead":   30,
    "body":   28,
    "body_s": 26,
    "small":  18,
    "huge":  150,
    "huge_s": 110,
}
MIN_BODY = 24    # 本文・箇条書き・リードの下限
MIN_CARD = 20    # カード内／階段図など、幅の狭い枠の下限
MIN_ANY  = 16    # スライド上のあらゆる文字の絶対下限（ページ番号14ptのみ例外）

# 重要な情報は 上端 t 〜 下端 (7.5 - b) の安全域に収める（会場の前列の頭・暗幕対策）
M = {"l": 0.95, "r": 0.95, "t": 0.55, "b": 1.05}
CONTENT_W = SLIDE_W - M["l"] - M["r"]

WARNINGS = []
ERRORS = []
SAFE_CHECK = [True]   # 全面ベタ塗りのスライドでは一時的に外す
CUR = [0]             # 検査メッセージ用の現在のスライド番号
SAFE_BOTTOM = SLIDE_H - 1.05

# ---------------------------------------------------------------- low level
def _rgb(hexs):
    return RGBColor.from_string(hexs)

def _set_run_font(run, size, bold=False, color="12253F", font=FONT, spc=None, italic=False):
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = _rgb(color)
    rPr = run._r.get_or_add_rPr()
    latin = rPr.find(qn("a:latin"))
    if latin is None:
        latin = etree.SubElement(rPr, qn("a:latin"))
        latin.set("typeface", font)
    prev = latin
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is not None:
            el.getparent().remove(el)
        el = etree.Element(qn(tag))
        el.set("typeface", font)
        prev.addnext(el)
        prev = el
    if spc is not None:
        rPr.set("spc", str(int(spc * 100)))

def _is_wide(ch):
    o = ord(ch)
    return (0x1100 <= o <= 0x115F or 0x2E80 <= o <= 0xA4CF or 0xAC00 <= o <= 0xD7A3
            or 0xF900 <= o <= 0xFAFF or 0xFE30 <= o <= 0xFE6F or 0xFF00 <= o <= 0xFF60
            or 0xFFE0 <= o <= 0xFFE6 or 0x3000 <= o <= 0x303F)

SAFE = 1.11   # フォント差（Meiryo/游ゴシック/Noto）を吸収する安全率

def text_w_in(text, size_pt):
    """文字列の実測に近い幅（インチ）。安全率込みで少し大きめに見積もる。"""
    em = size_pt / 72.0
    w = 0.0
    for ch in text:
        w += em * (1.0 if _is_wide(ch) else 0.55)
    return w * SAFE

def est_lines(text, box_w_in, size_pt):
    """折返し行数の見積り（明示改行も考慮）"""
    total = 0
    for para in text.split("\n"):
        if not para:
            total += 1
            continue
        line, n = 0.0, 1
        em = size_pt / 72.0 * SAFE
        for ch in para:
            cw = em * (1.0 if _is_wide(ch) else 0.55)
            if line + cw > box_w_in:
                n += 1
                line = cw
            else:
                line += cw
        total += n
    return total

def textbox(slide, x, y, w, h, text, size, bold=False, color="12253F",
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.35,
            space_after=0, spc=None, font=FONT, tag="", nowrap=False):
    if nowrap and align != PP_ALIGN.CENTER:
        nowrap = False          # 左右揃えの箱で折返しを切ると描画側で中央寄せに化ける
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = not nowrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        if space_after:
            p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = ln
        _set_run_font(r, size, bold, color, font, spc)
    if size < MIN_ANY and tag not in ("pn", "memo"):
        ERRORS.append(f"p{CUR[0]:02d} [小さすぎ] {tag or text[:14]!r} {size}pt < {MIN_ANY}pt")
    if tag in ("bullet", "lead") and size < MIN_BODY:
        ERRORS.append(f"p{CUR[0]:02d} [本文が下限割れ] {tag} {text[:16]!r} {size}pt < {MIN_BODY}pt")
    if SAFE_CHECK[0] and tag not in ("pn", "footnote") and y + h > SAFE_BOTTOM + 0.06:
        WARNINGS.append(f"p{CUR[0]:02d} [安全域外] {tag or text[:14]!r} 下端{y + h:.2f}in > {SAFE_BOTTOM:.2f}in")
    if nowrap:
        return tb
    need = est_lines(text, w, size) * (size / 72.0) * line_spacing
    if need > h + 0.04:
        WARNINGS.append(f"p{CUR[0]:02d} [overflow] {tag or text[:18]!r} 必要{need:.2f}in > 枠{h:.2f}in ({size}pt)")
    return tb

def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE,
         adj=None, shadow=False):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = _rgb(fill)
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = _rgb(line)
        s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    if not shadow:
        s.shadow.inherit = False
    if adj is not None:
        try:
            s.adjustments[0] = adj
        except Exception:
            pass
    s.text_frame.word_wrap = True
    return s

def bg(slide, color):
    r = rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=color)
    slide.shapes._spTree.remove(r._element)
    slide.shapes._spTree.insert(2, r._element)
    return r

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text or ""

def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

# ---------------------------------------------------------------- components
def header(slide, kicker, title, color_title="12253F", color_kicker="1160C4",
           rule="1160C4", y=None):
    """共通ヘッダ（ケッカー→タイトル→アクセントの罫）。戻り値＝本文開始 y"""
    y = M["t"] if y is None else y
    if kicker:
        textbox(slide, M["l"], y, CONTENT_W, 0.42, kicker.upper(), T["kicker"],
                True, color_kicker, spc=2.2, tag="kicker")
        y += 0.50
    size = T["title"] if text_w_in(title, T["title"]) <= CONTENT_W else T["title_s"]
    h = est_lines(title, CONTENT_W, size) * (size / 72.0) * 1.2
    textbox(slide, M["l"], y, CONTENT_W, h + 0.06, title, size, True, color_title,
            line_spacing=1.2, tag="title")
    y += h + 0.20
    if rule:
        rect(slide, M["l"], y, 1.55, 0.075, fill=rule)
        y += 0.075 + 0.34
    return y

def bullet_block(slide, x, y, w, points, size=None, gap=0.30, marker="1160C4",
                 color="12253F", max_h=None):
    """箇条書き。マーカーは図形で描き、行間をゆったり取る。"""
    size = size or T["body"]
    if len(points) > 4:
        ERRORS.append(f"p{CUR[0]:02d} [箇条書きが5項目以上] {len(points)}項目: {points[0][:16]!r}")
    for p_ in points:
        if len(p_) > 30:
            WARNINGS.append(f"p{CUR[0]:02d} [1行30字超] {len(p_)}字: {p_[:22]!r}")
    tw = w - 0.42
    # 収まらなければ 1 段だけ縮める（MIN_BODY 未満にはしない）
    while size > MIN_BODY and max_h:
        need = sum(est_lines(p, tw, size) * (size / 72.0) * 1.4 + gap for p in points) - gap
        if need <= max_h:
            break
        size -= 1
    if max_h:
        need = sum(est_lines(p, tw, size) * (size / 72.0) * 1.4 + gap for p in points) - gap
        while gap > 0.16 and need > max_h:
            gap -= 0.02
            need = sum(est_lines(p, tw, size) * (size / 72.0) * 1.4 + gap for p in points) - gap
        if need > max_h + 0.03:
            ERRORS.append(
                f"p{CUR[0]:02d} [箇条書きが枠に入らない] 必要{need:.2f}in > 枠{max_h:.2f}in "
                f"（{size}pt・{len(points)}項目）→ 文言を削るか項目を減らすこと")
    cy = y
    for p in points:
        n = est_lines(p, tw, size)
        h = n * (size / 72.0) * 1.4
        rect(slide, x + 0.02, cy + (size / 72.0) * 0.42, 0.17, 0.17, fill=marker,
             shape=MSO_SHAPE.OVAL)
        textbox(slide, x + 0.42, cy, tw, h + 0.05, p, size, False, color,
                line_spacing=1.4, tag="bullet")
        cy += h + gap
    return cy

def pill(slide, x, y, w, h, text, size, fill, color="FFFFFF", bold=True, min_size=14):
    while size > min_size and text_w_in(text, size) > w - 0.26:
        size -= 1
    s = rect(slide, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.5)
    tf = s.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    _set_run_font(r, size, bold, color)
    return s

def photo_slot(slide, x, y, w, h, caption="", hint="", tint="E7F0FB"):
    """写真プレースホルダ。写真を入れなくても『色面』として成立する見え方にする。"""
    rect(slide, x, y, w, h, fill=tint)
    rect(slide, x, y, 0.11, h, fill=C["brand"])
    if caption:
        cs = 22
        while cs > 18 and est_lines(caption, w - 0.68, cs) > 2:
            cs -= 1
        textbox(slide, x + 0.34, y + h - 1.05, w - 0.68, 0.85, caption, cs, True,
                C["ink"], anchor=MSO_ANCHOR.BOTTOM, line_spacing=1.3, tag="photocap")
    memo = "［写真を差し込む位置］" + (("　" + hint.replace("\n", " ")) if hint else "")
    hh = est_lines(memo, w - 0.68, 12) * (12 / 72.0) * 1.35
    textbox(slide, x + 0.34, y + 0.26, w - 0.68, hh + 0.06, memo, 12, False, "9BABBC",
            line_spacing=1.35, tag="memo")

def fn_height(text):
    """脚注が実際に占める高さ（余白込み）"""
    if not text:
        return 0.0
    return est_lines(text, CONTENT_W, T["small"]) * (T["small"] / 72.0) * 1.3 + 0.26


def footnote(slide, text):
    if not text:
        return
    h = est_lines(text, CONTENT_W, T["small"]) * (T["small"] / 72.0) * 1.3
    textbox(slide, M["l"], SLIDE_H - M["b"] - h, CONTENT_W, h + 0.05, text,
            T["small"], False, C["gray"], line_spacing=1.3, tag="footnote")

def pagenum(slide, n):
    textbox(slide, SLIDE_W - M["r"] - 0.9, SLIDE_H - 0.44, 0.9, 0.30, str(n), 14,
            False, "8895A5", align=PP_ALIGN.RIGHT, tag="pn")

def engage_strip(slide, text, y=None):
    """聴衆への投げかけを『帯』で見せる。単調さを壊す装置。"""
    if not text:
        return
    tw = CONTENT_W - 0.7
    size = 24
    while size > 20 and est_lines(text, tw, size) > 1:
        size -= 1
    n = est_lines(text, tw, size)
    h = max(0.86, n * (size / 72.0) * 1.3 + 0.32)
    y = SLIDE_H - M["b"] - h - 0.05 if y is None else y
    rect(slide, M["l"], y, CONTENT_W, h, fill=C["gold_lt"])
    rect(slide, M["l"], y, 0.11, h, fill=C["gold"])
    textbox(slide, M["l"] + 0.36, y + 0.14, tw, h - 0.26, text, size, True,
            C["gold_tx"], line_spacing=1.3, anchor=MSO_ANCHOR.MIDDLE, tag="engage")
